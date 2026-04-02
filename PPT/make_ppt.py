"""
生成周报PPT：DrugOperatorNet 进展汇报
基于模板 PPT/工作进展3.20.pptx
"""
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEMPLATE = "/home/data/jiangyun/cgi_data_pipeline5/PPT/工作进展3.20.pptx"
OUTPUT = "/home/data/jiangyun/cgi_data_pipeline5/PPT/工作进展4.3.pptx"

# Colors
WHITE = (255, 255, 255)
DARK_BLUE = (31, 73, 125)
LIGHT_BLUE = (173, 216, 230)
ORANGE = (204, 85, 0)
GRAY = (100, 100, 100)
BLACK = (0, 0, 0)
GREEN = (0, 128, 0)
RED = (180, 0, 0)

def add_tb(slide, l, t, w, h, text, fs=16, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fs)
    run.font.bold = bold
    if color:
        run.font.color.rgb = RGBColor(*color)
    return txBox

def add_tb_lines(slide, l, t, w, h, lines, fs=15, color=None, first_bold=False, first_color=None, first_fs=None):
    """Add textbox with multiple lines, each line is (text,) or (text, bold, color, fs)"""
    txBox = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            txt, lb, lc, lfs = (line + (None, None, None))[:4]
        else:
            txt, lb, lc, lfs = line, False, color, fs
        if i == 0 and first_bold:
            lb = True
        if i == 0 and first_color:
            lc = first_color
        if i == 0 and first_fs:
            lfs = first_fs
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(lfs or fs)
        run.font.bold = lb
        if lc:
            run.font.color.rgb = RGBColor(*lc)
    return txBox

def add_rect(slide, l, t, w, h, fill_color, line_color=None):
    from pptx.util import Cm
    shape = slide.shapes.add_shape(1, Cm(l), Cm(t), Cm(w), Cm(h))  # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_color)
    if line_color:
        shape.line.color.rgb = RGBColor(*line_color)
    else:
        shape.line.fill.background()
    return shape

# ============================================================
prs = Presentation(TEMPLATE)
W, H = prs.slide_width.cm, prs.slide_height.cm  # 33.9 x 19.1

def clear_text_shapes(slide, keep_bg=True):
    """Remove text boxes, keep pictures and connector lines if keep_bg"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    to_remove = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            to_remove.append(shape)
    for shape in to_remove:
        shape._element.getparent().remove(shape._element)

# ================================================================
# Slide 1: 封面
# ================================================================
s1 = prs.slides[0]
clear_text_shapes(s1)

add_tb(s1, 1, 4, 32, 4,
    "化学品-基因相互作用预测\nDrugOperatorNet: 基于扰动算子的端到端可解释深度学习",
    fs=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_tb(s1, 4, 11, 26, 2,
    "汇报人：\n日期：2026/04/03",
    fs=20, color=(220, 220, 220), align=PP_ALIGN.CENTER)

# ================================================================
# Slide 2: 研究背景与问题
# ================================================================
s2 = prs.slides[1]
clear_text_shapes(s2)

add_tb(s2, 0.5, 0.3, 20, 1.2, "研究背景与问题定义", fs=24, bold=True, color=DARK_BLUE)

# Left column: problem
add_tb_lines(s2, 0.5, 2.0, 15, 16, [
    "研究问题",
    "",
    "给定化学品 C 和基因 G，预测：",
    "C 是否能显著改变细胞系中 G 的表达？",
    "（二分类：上调/下调 = 1；无显著变化 = 0）",
    "",
    "数据集：L1000（LINCS）",
    "• 4 个细胞系：MCF7 / A375 / A549 / VCAP",
    "• 化学冷分割（Chemical Cold Split）：训练/测试无共享化合物",
    "• 5-fold Cross Validation",
    "• 标注：logFC > 1.5 且 q < 0.05 为正例",
    "",
    "核心挑战",
    "• 化学品泛化：固定指纹无法学习结构-活性关系",
    "• 机制可解释：预测需落地到具体药效团和调控通路",
    "• 类别不平衡：正例约占 15-25%",
], fs=14, first_bold=True, first_color=DARK_BLUE, first_fs=18)

# Right column: existing methods limitation
add_tb_lines(s2, 17, 2.0, 16, 16, [
    "现有方法的局限",
    "",
    "传统ML（Random Forest + Morgan FP）",
    "• 固定指纹，无法优化化学表示",
    "• 无法建模原子级别的结构-活性关系",
    "• 本文基线：AUC 0.789~0.887",
    "",
    "现有深度学习（DeepCE等）",
    "• 黑盒：药物嵌入 ⊕ 基因嵌入 → MLP",
    "• 无药效团对齐，可解释性差",
    "• 多数用预训练模型，非端到端",
    "",
    "本文方案",
    "• 算子范式：药物 = 基因空间上的扰动算子",
    "• T = I + UΣVᵀ，谱 Σ = 药物-基因耦合强度",
    "• GIN 端到端图学习 + 基因 k-mer 序列编码",
    "• 正交正则确保每个模式独立可解释",
], fs=14, first_bold=True, first_color=DARK_BLUE, first_fs=18)

# ================================================================
# Slide 3: 模型架构
# ================================================================
s3 = prs.slides[2]
clear_text_shapes(s3)

add_tb(s3, 0.5, 0.3, 32, 1.2, "DrugOperatorNet 模型架构", fs=24, bold=True, color=DARK_BLUE)

# Architecture description as ASCII-style boxes
arch_text = [
    "输入层",
    "药物分子图 G_mol（原子节点 + 键边）    基因序列 DNA/RNA（k-mer tokenize）",
    "",
    "编码层",
    "GIN（3层消息传递）                      CNN-MultiHead（k-mer → 局部模式）",
    "└─ 原子嵌入 H_atom [B, N_a, H]         └─ 基因全局表示 V_g [B, H]",
    "└─ 药效团提取器（PharmacophoreExtractor）",
    "   r 个可学习 query slots × 原子 cross-attention",
    "   → 药效团模式 U [B, r, H]（r个独立药效团）",
    "",
    "算子层（PerturbationOperator）",
    "• 基因多头读取器：V_g → r 个基因模式 [B, r, H]",
    "• 模式对齐耦合：spectrum_k = σ(v_k · h_g_mode_k)  [B, r]",
    "• delta_h = Σ_k spectrum_k * u_k  [B, H]",
    "  （药物对基因空间的有效扰动）",
    "",
    "分类层",
    "• 特征拼接：[V_g_global, delta_h]  [B, 2H]",
    "• MLP(2H → H → 1) + Sigmoid",
    "",
    "正则化",
    "• 正交正则 L_ortho：Gram(U) ≈ I，确保药效团模式独立",
    "• lam_ortho = 0.1（消融验证有效）",
]
add_tb_lines(s3, 0.5, 1.8, 32, 16, arch_text, fs=13, first_bold=True, first_color=DARK_BLUE, first_fs=15)

# ================================================================
# Slide 4: 实验结果 - 主要性能对比
# ================================================================
s4 = prs.slides[3]
clear_text_shapes(s4)

add_tb(s4, 0.5, 0.3, 32, 1.2, "实验结果：4细胞系性能对比", fs=24, bold=True, color=DARK_BLUE)

# Table header
add_tb(s4, 0.5, 1.8, 32, 0.8,
    "| 细胞系     | Morgan FP | MoE      | MoE+Target | DrugOperatorNet | GIN净增益 |",
    fs=13, bold=True, color=DARK_BLUE)
add_tb_lines(s4, 0.5, 2.5, 32, 10, [
    "|------------|-----------|----------|------------|-----------------|-----------|",
    "| MCF7       | 0.8710    | 0.8935   | 0.8943     | 0.8923          | +0.0233   |",
    "| A375       | 0.8870    | 0.9035   | 0.9040     | (运行中)        | +0.0170   |",
    "| A549       | 0.8432    | 0.8703   | 0.8718     | (运行中)        | +0.0286   |",
    "| VCAP       | 0.7894    | 0.8132   | 0.8135     | (运行中)        | +0.0241   |",
    "",
    "关键结论：",
    "1. GIN端到端图学习 在4/4细胞系上一致优于固定Morgan指纹，净增益 +0.017~+0.029",
    "2. DrugOperatorNet（0.8923）≈ MoE（0.8935），用更少参数（918K vs 1.06M）达到相同效果",
    "3. 化学冷分割场景下GIN的结构归纳偏置优势更明显（固定指纹对未见化合物泛化弱）",
    "",
    "正在运行：A375/A549/VCAP 的 DrugOperatorNet（cuda:0/2/3）",
], fs=13)

add_tb(s4, 0.5, 13.0, 32, 5.5,
    "MCF7 消融实验（验证各组件贡献）", fs=16, bold=True, color=DARK_BLUE)
add_tb_lines(s4, 0.5, 14.0, 32, 4.5, [
    "| 变体            | AUC    | 说明                          |",
    "|-----------------|--------|-------------------------------|",
    "| full（谱路由）  | 0.8868 | 最差：MoE路由冲突             |",
    "| no_spectrum     | 0.8917 | 传统MoE路由，无算子谱         |",
    "| no_moe（主模型）| 0.8923 | 纯算子，MLP分类头，最优       |",
    "| no_moe_noortho  | 0.8915 | 去正交正则，-0.0008 AUC       |",
], fs=13)

# ================================================================
# Slide 5: 消融分析与结论
# ================================================================
s5 = prs.slides[4]
clear_text_shapes(s5)

add_tb(s5, 0.5, 0.3, 32, 1.2, "消融实验结论 & 算子范式有效性", fs=24, bold=True, color=DARK_BLUE)

add_tb_lines(s5, 0.5, 2.0, 15.5, 16, [
    "四个核心结论",
    "",
    "1. MoE融合无效",
    "   delta_h已编码交互信息，MoE的LB Loss",
    "   引入梯度冲突，反降 0.8923 → 0.8868",
    "",
    "2. 纯算子（no_moe）最优",
    "   AUC=0.8923，918K参数（最少）",
    "   验证了算子范式的有效性",
    "",
    "3. 正交正则有效",
    "   lam_ortho 0.1 vs 0.0：+0.0008 AUC",
    "   主要价值：保证药效团模式的独立可解释性",
    "",
    "4. 谱信号适合分类而非路由",
    "   spectrum作为特征 > spectrum作为路由信号",
    "",
    "确定主模型：DrugOperatorNet (no_moe)",
    "= train_operator_moe.py --ablation no_moe",
], fs=14, first_bold=True, first_color=DARK_BLUE, first_fs=18)

add_tb_lines(s5, 17.5, 2.0, 15.5, 16, [
    "待做实验（NMI必要）",
    "",
    "统计显著性",
    "□ 5-fold CV on all 4 cell lines",
    "□ mean ± std AUC，t-test vs baseline",
    "",
    "SOTA对比",
    "□ DeepCE（Song et al. 2021）",
    "□ DECIPHIR 或同期方法",
    "□ 需实现或找开源代码运行",
    "",
    "可解释性分析（详见下页）",
    "□ 谱分析：正/负样本区分",
    "□ 药效团热图：原子注意力",
    "□ 基因注意力：GO富集",
    "",
    "当前进度：~60% for NMI",
], fs=14, first_bold=True, first_color=DARK_BLUE, first_fs=18)

# ================================================================
# Slide 6: 可解释性工作计划
# ================================================================
s6 = prs.slides[5]
clear_text_shapes(s6)

add_tb(s6, 0.5, 0.3, 32, 1.2, "可解释性分析工作计划（NMI核心要求）", fs=24, bold=True, color=DARK_BLUE)

add_tb_lines(s6, 0.5, 2.0, 10, 16, [
    "① 交互谱分析",
    "（analyze_spectrum.py）",
    "",
    "• 对4细胞系正/负样本各提取谱",
    "• 可视化：每个模式k的激活分布",
    "• 期望：正样本谱值显著更高",
    "• 输出：violinplot / heatmap",
    "",
    "• 化合物聚类：按谱向量t-SNE",
    "  相同靶标的化合物应聚在一起",
    "",
    "• 工具函数已有：save_spectrum=True",
    "  spectrum.npy保存于results/",
], fs=13, first_bold=True, first_color=DARK_BLUE, first_fs=16)

add_tb_lines(s6, 11.5, 2.0, 10, 16, [
    "② 药效团热图",
    "（分子注意力可视化）",
    "",
    "• 保存PharmacophoreExtractor",
    "  的原子注意力权重",
    "• 在RDKit中绘制分子结构，",
    "  按注意力权重着色",
    "• 期望：已知药效团位置高亮",
    "",
    "• 每个模式k对应一类药效团",
    "  （如：模式1=芳香环，模式2=氢键供体）",
    "",
    "• 代码修改：train_drug_operator_v2.py",
    "  保存atom_attn_weights.npy",
], fs=13, first_bold=True, first_color=DARK_BLUE, first_fs=16)

add_tb_lines(s6, 23, 2.0, 10, 16, [
    "③ 基因注意力分析",
    "（GeneMultiHeadReader）",
    "",
    "• 保存基因序列的注意力权重",
    "• 映射到基因组坐标系",
    "• 与已知调控元件对比：",
    "  - 转录因子结合位点（JASPAR）",
    "  - 增强子/启动子区域",
    "",
    "• GO富集分析：",
    "  Top-k高注意力基因 → GO BP",
    "• 期望：细胞系特异性通路富集",
    "  MCF7：ER信号/MAPK",
    "  A375：BRAF/MEK/ERK",
    "",
    "• 工具：save_gene_attn=True",
], fs=13, first_bold=True, first_color=DARK_BLUE, first_fs=16)

# ================================================================
# Slide 7: 论文写作计划
# ================================================================
s7 = prs.slides[6]
clear_text_shapes(s7)

add_tb(s7, 0.5, 0.3, 32, 1.2, "NMI论文写作计划", fs=24, bold=True, color=DARK_BLUE)

add_tb_lines(s7, 0.5, 2.0, 15.5, 16, [
    "论文结构（目标期刊：NMI）",
    "",
    "Abstract（250字）",
    "• 问题 → 现有不足 → 本文方案",
    "• DrugOperatorNet + 4细胞系 + 可解释性",
    "",
    "Introduction（~1500字）",
    "• CGI预测在毒理学/药物发现中的意义",
    "• 化学品-基因关系的复杂性",
    "• 现有方法：ML基线 / 黑盒DL / 预训练模型",
    "• 本文贡献：算子范式 / 端到端 / 可解释",
    "",
    "Methods（~2000字）",
    "• 数据集描述（L1000/化学冷分割）",
    "• DrugOperatorNet架构",
    "  - GIN药物编码",
    "  - k-mer基因序列编码",
    "  - PerturbationOperator（T=I+UΣVᵀ）",
    "  - 正交正则损失",
    "• 训练细节",
], fs=13, first_bold=True, first_color=DARK_BLUE, first_fs=17)

add_tb_lines(s7, 17.5, 2.0, 15.5, 16, [
    "论文结构（续）",
    "",
    "Results（~2500字）",
    "• R1: 主性能对比（4细胞系 vs SOTA）",
    "• R2: 消融实验（算子各组件贡献）",
    "• R3: GIN vs Morgan FP（端到端必要性）",
    "• R4: 可解释性分析",
    "  - 谱分析（正/负样本区分）",
    "  - 药效团热图（分子级别）",
    "  - 基因注意力GO富集",
    "",
    "Discussion（~1000字）",
    "• 算子范式的生物学意义",
    "• 模式独立性 → 多机制并行",
    "• 化学结构归纳偏置的优势",
    "• 局限性：序列→结构层次/动态调控",
    "",
    "当前写作建议",
    "□ 先写Methods（已稳定）",
    "□ 等全部5-fold结果后写Results",
    "□ Introduction最后写",
    "□ 预计投稿时间：2026年6-7月",
], fs=13, first_bold=True, first_color=DARK_BLUE, first_fs=17)

# ================================================================
# Slide 8: 下一步计划
# ================================================================
s8 = prs.slides[7]
# Clear all
clear_text_shapes(s8)

add_tb(s8, 0.5, 0.3, 32, 1.2, "下一步计划（优先级排序）", fs=24, bold=True, color=DARK_BLUE)

add_tb_lines(s8, 0.5, 2.0, 32, 16, [
    "P0（本周完成）",
    "□ 等待 A375/A549/VCAP DrugOperatorNet 结果（cuda:0/2/3）",
    "□ 运行 analyze_spectrum.py：4细胞系谱可视化",
    "□ 修改代码保存原子注意力权重，生成MCF7样例热图",
    "",
    "P1（下周）",
    "□ 5-fold CV：DrugOperatorNet 在4细胞系全部5个fold",
    "□ 调研并实现SOTA对比方法（DeepCE / DECIPHIR）",
    "□ 基因注意力分析 → GO富集（MCF7/A375各一个细胞系）",
    "",
    "P2（两周内）",
    "□ 完成论文Methods部分初稿",
    "□ 统计显著性检验（paired t-test，5 seeds × 4 cell lines）",
    "□ 准备投稿图（高清PDF，Nature格式）",
    "",
    "P3（一个月）",
    "□ 完成论文全文初稿",
    "□ 预印本上传 bioRxiv",
    "□ 提交 NMI",
    "",
    "当前进度评估：实验部分约60%，可解释性约20%，论文约5%",
], fs=15, first_bold=True, first_color=DARK_BLUE, first_fs=18)

# Save
prs.save(OUTPUT)
print(f"PPT saved to: {OUTPUT}")
